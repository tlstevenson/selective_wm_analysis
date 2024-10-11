# -*- coding: utf-8 -*-
"""
Created on Thu Dec 21 14:58:30 2023

@author: tanne
"""
import init

import pyutils.utils as utils
from sys_neuro_tools import plot_utils, fp_utils
import numpy as np
import pandas as pd
from enum import StrEnum
import matplotlib.pyplot as plt
import scipy.signal as sig
from scipy.optimize import curve_fit
import os.path as path
import os
from glob import glob
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import time
import copy
import warnings

# %% trial start timestamps

# def get_trial_start_ts(fp_start_ts, sess_data):
#     fp_ts_diffs = np.diff(fp_start_ts)
#     beh_ts_diffs = np.diff(trial_data['trialtime'])



# %% Session Limiting Lookup Table

# Some sessions have poor signal quality due to patch cord connection issues, so this function is used to exclude fp signal and behavioral data from analyses
# for such sessions, specify time ranges to include in seconds
__sess_exclusions = {96556: {'PL': [0,3871]},
                     101853: {'PL': [0,832]},
                     101906: {'PL': [0,1072], 'DMS': [0,2188]},
                     102186: {'PL': [[0,349], [350,1901]]},
                     102235: {'PL': [[0,57], [57.5,3659], [3659.5]]}, #, 'DMS': [[0,3659], [3659.5]]},
                     101729: {'DMS': [[0,1890], [1890.5,3398]]},
                     101926: {'PL': [0,542.5]},
                     101872: {'PL': [[0,113.5], [114.5,1026.5]]},
                     102053: {'PL': [0,2571], 'DMS': [0,2571]},
                     102201: {'PL': [0,3442]},
                     102340: {'PL': [[0,60.5], [61,1286], [1286.5]]},
                     102394: {'PL': [0,3674], 'DMS': [0,3126.5]},
                     101717: {'PL': [0,1993]},
                     102620: {'PL': [[0,411], [850,1337.5]], 'DMS': [0,849.5]},
                     102530: {'PL': [0,925], 'DMS': [0,1692]},
                     102580: {'PL': [0,2144], 'DMS': [0,1110.5]},
                     102614: {'PL': [[0,2072.5], [2073.5]]},
                     102524: {'PL': [0,3718.5]},
                     102605: {'PL': [[0,880.5], [881.5]]},
                     101896: {'PL': [[0,442.5], [443.5,1220], [1220.5]]},
                     102318: {'DMS': [[0,1590.5], [1591]]},
                     102367: {'PL': [0,382], 'DMS': [0,382]}
                    }

__sess_ignore = [100323, 101581]

# choose 405 over 420 when there are sessions with both for 3.6
preferred_isos = {182: ['405', '420'], 202: ['405', '420'], 179: ['420', '405'],
                  180: ['420', '405'], 188: ['420', '405'], 191: ['420', '405'], 207: ['420', '405']}

# %% Loading methods

def load_fp_data(loc_db, sess_ids, isos=None, ligs=None, reload=False, fit_baseline=True, lig_lpf=10, iso_lpf=1):
    ''' Loads and processes fiber photometry data, optionally choosing which wavelengths are ligand and isosbsestic.'''

    start = time.perf_counter()

    # get fiber photometry data
    fp_data = loc_db.get_sess_fp_data(utils.flatten(sess_ids), reload=reload)

    print('Retreived FP data for {} session(s) in {:.1f} s'.format(len(utils.flatten(sess_ids)), time.perf_counter()-start))

    # separate into different dictionaries
    implant_info = fp_data['implant_info']
    fp_data = fp_data['fp_data']

    if ligs is None:
        ligs = ['490', '465']
    elif not utils.is_list(ligs):
        ligs = [ligs]

    ligs = np.array(ligs)

    # process the signals
    for subj_id in sess_ids.keys():
        
        if isos is None:
            subj_isos = preferred_isos[subj_id]
        elif not utils.is_list(isos):
            subj_isos = [isos]
            
        subj_isos = np.array(subj_isos)
        
        for sess_id in sess_ids[subj_id]:

            fp_data[subj_id][sess_id]['processed_signals'] = {}

            raw_signals = fp_data[subj_id][sess_id]['raw_signals']

            for region in raw_signals.keys():
                lig_sel = np.array([k in raw_signals[region].keys() for k in ligs])
                iso_sel = np.array([k in raw_signals[region].keys() for k in subj_isos])

                if sum(lig_sel) > 1:
                    lig = ligs[0]
                    print('Found {} matching ligand wavelengths. Choosing {} nm'.format(sum(lig_sel), ligs[0]))
                elif sum(lig_sel) == 1:
                    lig = ligs[lig_sel][0]
                else:
                    raise Exception('No ligand wavelength found')

                if sum(iso_sel) > 1:
                    iso = subj_isos[0]
                    print('Found {} matching iso wavelengths. Choosing {} nm'.format(sum(iso_sel), subj_isos[0]))
                elif sum(iso_sel) == 1:
                    iso = subj_isos[iso_sel][0]
                else:
                    raise Exception('No isosbestic wavelength found')

                raw_lig = raw_signals[region][lig]
                raw_iso = raw_signals[region][iso]

                start = time.perf_counter()

                fp_data[subj_id][sess_id]['processed_signals'][region] = get_all_processed_signals(raw_lig, raw_iso,
                                                                                                   t=fp_data[subj_id][sess_id]['time'],
                                                                                                   sess_id=sess_id, region=region,
                                                                                                   trial_start_ts=fp_data[subj_id][sess_id]['trial_start_ts'],
                                                                                                   fit_baseline=fit_baseline,
                                                                                                   lig_lpf=lig_lpf, iso_lpf=iso_lpf)

                print('  Processed FP data in {:.1f} s'.format(time.perf_counter()-start))


    return fp_data, implant_info


def get_all_processed_signals(raw_lig, raw_iso, t, sess_id=None, region=None, trial_start_ts=None, fit_baseline=True,
                              lig_lpf=10, iso_lpf=1, filter_dropout_outliers=True):
    ''' Gets all possible processed signals and intermediaries for the given raw signals.
        Will check to see if any signals should be excluded. Also will also optionally exclude signals before and after the behavior.'''

    # initialize signal variables
    empty_signal = np.full_like(raw_lig, np.nan)
    filtered_lig = empty_signal.copy()
    filtered_iso = empty_signal.copy()
    baseline_lig = empty_signal.copy()
    baseline_iso = empty_signal.copy()
    baseline_corr_lig = empty_signal.copy()
    baseline_corr_iso = empty_signal.copy()
    fitted_iso = empty_signal.copy()
    fitted_baseline_corr_iso = empty_signal.copy()
    dff_iso = empty_signal.copy()
    dff_iso_baseline = empty_signal.copy()

    # if we are to ignore this session, just leave all processed signals as nan
    if sess_id is None or not sess_id in __sess_ignore:

        sr = 1/np.mean(np.diff(t))

        if not trial_start_ts is None and not t is None:
            start_ts_idx = np.argmin(np.abs(t - trial_start_ts[0]))
            # end the signal a little past the last trial start timestamp (which is sent before the trial interrupted when ending the recording)
            end_ts_idx = np.argmin(np.abs(t - trial_start_ts[-1] - 20))
        else:
            start_ts_idx = 0
            end_ts_idx = len(raw_lig)

        # In order to exclude parts of the signal, need to pass in the time, session id, and region
        if not sess_id is None and not region is None and sess_id in __sess_exclusions.keys() and region in __sess_exclusions[sess_id].keys():

            t_ranges = __sess_exclusions[sess_id][region]

            if not utils.is_list(t_ranges[0]):
                t_ranges = [t_ranges]

        else:
            t_ranges = [[t[0], t[-1]]]

        # process each time range independently
        for t_range in t_ranges:
            start_idx = np.argmin(np.abs(t - t_range[0]))
            if len(t_range) == 1:
                # go to the end of the signal
                end_idx = len(t)
            else:
                end_idx = np.argmin(np.abs(t - t_range[1]))

            if start_idx < start_ts_idx:
                start_idx = start_ts_idx

            if end_idx > end_ts_idx:
                end_idx = end_ts_idx

            sub_raw_lig = raw_lig[start_idx:end_idx]
            sub_raw_iso = raw_iso[start_idx:end_idx]

            if filter_dropout_outliers:
                # since dropouts have a high frequency component, use a highpass filter to find any outliers
                hpf_lig = utils.z_score(fp_utils.filter_signal(sub_raw_lig, 50, sr, filter_type='highpass'))
                hpf_iso = utils.z_score(fp_utils.filter_signal(sub_raw_iso, 50, sr, filter_type='highpass'))
                lig_sel = np.abs(hpf_lig) >= 15
                iso_sel = np.abs(hpf_iso) >= 15
                drop_sel = lig_sel | iso_sel

                filtered = False
                outlier_filtered_lig = sub_raw_lig.copy()
                outlier_filtered_iso = sub_raw_iso.copy()
                if sum(drop_sel) > 1:
                    # find start and end of each dropout and set all intervening values to nan
                    drop_idxs = np.flatnonzero(drop_sel)
                    first_idx = drop_idxs[0]
                    while first_idx < drop_idxs[-1]:
                        # find end of dropout as last threshold cross within time interval
                        last_idx = drop_idxs[(drop_idxs >= first_idx) & (drop_idxs < first_idx+5*sr)][-1]

                        outlier_filtered_lig[first_idx:last_idx+1] = np.nan
                        outlier_filtered_iso[first_idx:last_idx+1] = np.nan
                        remaining_idxs = drop_idxs[drop_idxs > last_idx]
                        if len(remaining_idxs) > 0:
                            first_idx = remaining_idxs[0]
                        else:
                            break

                    filtered = True

                if filtered:
                    print('Session {} had dropout outliers in {}'.format(sess_id, region))
                    # fig, axs = plt.subplots(2,1, sharex=True)
                    # axs[0].plot(t[start_idx:end_idx], sub_raw_lig)
                    # axs[0].plot(t[start_idx:end_idx], outlier_filtered_lig)
                    # axs[1].plot(t[start_idx:end_idx], sub_raw_iso)
                    # axs[1].plot(t[start_idx:end_idx], outlier_filtered_iso)
                    # fig.suptitle('Session {} - {}'.format(sess_id, region))

                    # e=1

                sub_raw_lig = outlier_filtered_lig
                sub_raw_iso = outlier_filtered_iso


            sub_filt_lig = fp_utils.filter_signal(sub_raw_lig, lig_lpf, sr)
            sub_filt_iso = fp_utils.filter_signal(sub_raw_iso, iso_lpf, sr)

            # isosbestic correction
            sub_dff_iso, sub_fitted_iso = fp_utils.calc_iso_dff(sub_filt_lig, sub_filt_iso, t[start_idx:end_idx])

            sub_empty_signal = np.full_like(sub_raw_lig, np.nan)
            sub_baseline_lig = sub_empty_signal.copy()
            sub_baseline_iso = sub_empty_signal.copy()
            sub_baseline_corr_lig = sub_empty_signal.copy()
            sub_baseline_corr_iso = sub_empty_signal.copy()
            sub_fitted_baseline_iso = sub_empty_signal.copy()
            sub_dff_iso_baseline = sub_empty_signal.copy()

            if fit_baseline:

                # Try hybrid approach with baseline correction to approximate photobleaching before calculating dF/F
                try:
                    sub_baseline_lig = fp_utils.fit_baseline(sub_filt_lig)
                    sub_baseline_iso = fp_utils.fit_baseline(sub_filt_iso)

                    # first subtract the baseline fit to each signal to correct for photobleaching
                    sub_baseline_corr_lig = sub_filt_lig - sub_baseline_lig
                    sub_baseline_corr_iso = sub_filt_iso - sub_baseline_iso

                    # scale the isosbestic signal to best fit the ligand-dependent signal
                    sub_fitted_baseline_iso = fp_utils.fit_signal(sub_baseline_corr_iso, sub_baseline_corr_lig, t[start_idx:end_idx])

                    # then use the baseline corrected signals to calculate dF, which is a residual fluorescence
                    sub_dff_iso_baseline = ((sub_baseline_corr_lig - sub_fitted_baseline_iso)/sub_baseline_lig)*100

                except RuntimeError as error:
                    print(str(error))

            filtered_lig[start_idx:end_idx] = sub_filt_lig
            filtered_iso[start_idx:end_idx] = sub_filt_iso
            baseline_lig[start_idx:end_idx] = sub_baseline_lig
            baseline_iso[start_idx:end_idx] = sub_baseline_iso
            baseline_corr_lig[start_idx:end_idx] = sub_baseline_corr_lig
            baseline_corr_iso[start_idx:end_idx] = sub_baseline_corr_iso
            fitted_iso[start_idx:end_idx] = sub_fitted_iso
            fitted_baseline_corr_iso[start_idx:end_idx] = sub_fitted_baseline_iso
            dff_iso[start_idx:end_idx] = sub_dff_iso
            dff_iso_baseline[start_idx:end_idx] = sub_dff_iso_baseline

    return {'raw_lig': raw_lig,
            'raw_iso': raw_iso,
            'filtered_lig': filtered_lig,
            'filtered_iso': filtered_iso,
            'baseline_lig': baseline_lig,
            'baseline_iso': baseline_iso,
            'baseline_corr_lig': baseline_corr_lig,
            'baseline_corr_iso': baseline_corr_iso,
            'fitted_iso': fitted_iso,
            'fitted_baseline_corr_iso': fitted_baseline_corr_iso,
            'dff_iso': dff_iso,
            'z_dff_iso': utils.z_score(dff_iso),
            'dff_iso_baseline': dff_iso_baseline,
            'z_dff_iso_baseline': utils.z_score(dff_iso_baseline)}

# %% Analysis Methods
exp_decay_form = lambda t, a, tau, b: a*np.exp(-t/tau) + b
exp_decay_bounds = ([     0,      0, -np.inf],
                    [np.inf, np.inf,  np.inf])

def calc_peak_properties(signal, t, filter_params={}, peak_find_params={}, fit_decay=True):
    if len(signal) == 0:
        return {}
           
    dt = np.nanmean(np.diff(t))
    
    # read params
    filt = filter_params.get('filter', True)
    use_filt_signal = filter_params.get('use_filt_signal_props', False)
    cutoff_f = filter_params.get('cutoff_f', 10)
    peak_tmax = peak_find_params.get('peak_tmax', np.median(t))
    tau_tmax = peak_find_params.get('tau_tmax', t[-1])
    min_dist = peak_find_params.get('min_dist', 0)
    peak_edge_buffer = peak_find_params.get('peak_edge_buffer', 2*dt)
    peak_edge_buffer = int(np.ceil(peak_edge_buffer/dt))

    # smooth signal by filtering
    if filt:
        filt_signal = fp_utils.filter_signal(signal, cutoff_f, 1/dt)
    else:
        filt_signal = signal

    # fill nans so peak finding doesn't have issues
    filt_signal = fp_utils.fill_signal_nans(filt_signal)
    
    if use_filt_signal:
        peak_prop_signal = filt_signal
    else:
        peak_prop_signal = signal

    t0_idx = np.argmin(np.abs(t))

    peak_tmax_idx = np.argmin(np.abs(t - peak_tmax))
    tau_tmax_idx = np.argmin(np.abs(t - tau_tmax))
    
    if ((fit_decay and np.sum(np.isnan(signal[t0_idx:tau_tmax_idx])) > 0.25*(tau_tmax_idx - t0_idx)) or
       (not fit_decay and np.sum(np.isnan(signal[t0_idx:peak_tmax_idx])) > 0.25*(peak_tmax_idx - t0_idx))):
        return {}

    min_peak_dist_idxs = int(min_dist/dt)

    all_peaks, _ = sig.find_peaks(filt_signal)

    peak_idx = None
    if len(all_peaks) > 0:
        peaks = all_peaks[(all_peaks > t0_idx) & (all_peaks < peak_tmax_idx)]
        if len(peaks) > 0:
            # prioritize peaks in the middle of the range
            middle_peaks = peaks[(peaks > t0_idx + peak_edge_buffer) & (peaks < peak_tmax_idx - peak_edge_buffer)]
            if len(middle_peaks) > 0:
                peak_idx = middle_peaks[np.argmax(peak_prop_signal[middle_peaks])]
            else:
                peak_idx = peaks[np.argmax(peak_prop_signal[peaks])]

    if peak_idx is None:
        return {}

    valleys, _ = sig.find_peaks(-filt_signal)

    peak_right_idx = None
    peak_left_idx = None
    if len(valleys) > 0:
        left_valleys = valleys[valleys < peak_idx-min_peak_dist_idxs]
        right_valleys = valleys[(valleys > peak_idx+min_peak_dist_idxs) & (valleys < tau_tmax_idx)]

        if len(left_valleys) > 0:
            peak_left_idx = left_valleys[-1]
            
            # if there is a sub-peak between the max peak and t=0 that would shorten the calculated peak magnitude, 
            # move the peak edge to the next valley until that valley is lower than t=0
            i = -2
            while peak_prop_signal[t0_idx] < peak_prop_signal[peak_left_idx] and len(left_valleys) >= abs(i):
                peak_left_idx = left_valleys[i]
                i -= 1

        if len(right_valleys) > 0:
            peak_right_idx = right_valleys[0]

    if peak_left_idx is None:
        peak_left_idx = t0_idx

    if peak_right_idx is None:
        peak_right_idx = np.argmin(peak_prop_signal[peak_idx:tau_tmax_idx]) + peak_idx

    if peak_right_idx >= tau_tmax_idx:
        peak_right_idx = tau_tmax_idx
        lowest_side_idx = peak_left_idx
    else:
        idxs = np.array([peak_left_idx, peak_right_idx])
        lowest_side_idx = idxs[np.argmin(peak_prop_signal[idxs])]
        
    peak_height = peak_prop_signal[peak_idx] - peak_prop_signal[lowest_side_idx]
    peak_time = t[peak_idx]

    # manually find peak width
    peak_width_y = peak_prop_signal[peak_idx] - peak_height*0.5
    y_sel = np.flatnonzero(peak_prop_signal < peak_width_y)
    left_idxs = y_sel[(y_sel < peak_idx) & (y_sel > peak_left_idx)]
    right_idxs = y_sel[(y_sel > peak_idx) & (y_sel < peak_right_idx)]

    if len(left_idxs) > 0:
        left_idx = left_idxs[-1]
    else:
        left_idx = peak_left_idx

    if len(right_idxs) > 0:
        right_idx = right_idxs[0]
    else:
        right_idx = peak_right_idx

    peak_width_left = t[left_idx]
    peak_width_right = t[right_idx]
    peak_width = peak_width_right - peak_width_left

    #     # get peak height and width with peak finding routines
    #     # need to pass in all peak indexes for algorithm to work properly
    #     all_peaks_peak_idx = np.argmin(np.abs(all_peaks - peak_idx))
    #     peak_heights, peak_left_idxs, peak_right_idxs = sig.peak_prominences(avg_signal, all_peaks)
    #     peak_left_idx = peak_left_idxs[all_peaks_peak_idx]
    #     peak_right_idx = peak_right_idxs[all_peaks_peak_idx]
    #     if peak_right_idx > tau_tmax_idx:
    #         peak_right_idx = tau_tmax_idx
    #         lowest_side_idx = peak_left_idx
    #     else:
    #         # overwrite peak height with the height from the lowest peak edge
    #         idxs = np.array([peak_left_idx, peak_right_idx])
    #         lowest_side_idx = idxs[np.argmin(avg_signal[idxs])]

    #     peak_heights[all_peaks_peak_idx] = avg_signal[peak_idx] - avg_signal[lowest_side_idx]
    #     peak_widths, peak_widths_y, peak_widths_left, peak_widths_right = sig.peak_widths(avg_signal, all_peaks, rel_height=0.5,
    #                                                                                      prominence_data=(peak_heights, peak_left_idxs, peak_right_idxs))
    #     peak_height = peak_heights[all_peaks_peak_idx]
    #     peak_width = peak_widths[all_peaks_peak_idx]*dt
    #     peak_width_y = peak_widths_y[all_peaks_peak_idx]
    #     peak_width_left = peak_widths_left[all_peaks_peak_idx]*dt + t[0]
    #     peak_width_right = peak_widths_right[all_peaks_peak_idx]*dt + t[0]


    # if np.random.random() < 0.01: # or peak_height < 0
        # _, ax = plt.subplots(1,1)
        # ax.plot(t, signal)
        # ax.plot(t, filt_signal)
        # plot_utils.plot_dashlines([0, peak_tmax], ax=ax)
        # ax.plot(t[peak_idx], peak_prop_signal[peak_idx], marker=7, markersize=10, color='C1')
        # ax.vlines(t[peak_idx], peak_prop_signal[peak_idx]-peak_height, peak_prop_signal[peak_idx], color='C2', linestyles='dashed')
        # ax.plot([t[peak_left_idx], t[peak_right_idx]], [peak_prop_signal[peak_left_idx], peak_prop_signal[peak_right_idx]], color='C2')
        # ax.hlines(peak_width_y, peak_width_left, peak_width_right, color='C3')

    # fit exponential decay to falling peak
    if fit_decay:
        # add a few extra data points to help the fit be an actual decay
        decay_end_idx = int(peak_idx + 1.2*(peak_right_idx - peak_idx))
        sub_signal = peak_prop_signal[peak_idx:decay_end_idx]
        # fix baseline to be the bottom of the peak
        sub_bounds = copy.copy(exp_decay_bounds)
        decay_baseline = np.min(sub_signal)
        sub_bounds[0][-1] = decay_baseline-0.01 # this needs to be slightly smaller than the upper bound
        sub_bounds[1][-1] = decay_baseline
        decay_params = curve_fit(exp_decay_form, t[peak_idx:decay_end_idx] - peak_time, sub_signal, bounds=sub_bounds)[0]
        decay_form = exp_decay_form
        decay_tau = decay_params[1]
    else:
        decay_tau = np.nan
        decay_form = lambda x: np.nan
        decay_params = []

    return {'peak_time': peak_time, 'peak_height': peak_height, 'peak_width': peak_width,
            'peak_width_info': {'y': peak_width_y, 't_lims': [peak_width_left, peak_width_right]},
            'peak_end_time': t[peak_right_idx], 'decay_tau': decay_tau, 'decay_params': decay_params, 'decay_form': decay_form}


def calc_iqr_multiple(table, group_by_cols, parameters):
    table = table.copy()
    if not utils.is_list(group_by_cols):
        group_by_cols = [group_by_cols]

    if not utils.is_list(parameters):
        parameters = [parameters]

    iqr_mult_keys = {param: 'iqr_mult_' + param for param in parameters}
    # initialize iqr columns
    for param in parameters:
        if not iqr_mult_keys[param] in table:
            table[iqr_mult_keys[param]] = np.nan

    groupings = table[group_by_cols].drop_duplicates().to_numpy()
    for group in groupings:
        sel = np.all(np.array([(table[group_col] == group_val) for group_col, group_val in zip(group_by_cols, group)]), axis=0)

        sub_table = table.loc[sel]
        for param in parameters:
            # with warnings.catch_warnings():
            #     warnings.simplefilter('error')
            #     try:
            q1, q3 = np.nanquantile(sub_table[param], [0.25, 0.75])
                # except:
                #     e=1
            iqr = q3 - q1
            below_sel = sub_table[param] < q1
            above_sel = sub_table[param] > q3
            sub_table.loc[below_sel, iqr_mult_keys[param]] = (sub_table.loc[below_sel, param] - q1)/iqr
            sub_table.loc[above_sel, iqr_mult_keys[param]] = (sub_table.loc[above_sel, param] - q3)/iqr

        table.loc[sel, iqr_mult_keys.values()] = sub_table[iqr_mult_keys.values()]

    return table

# %% Plotting Methods

def view_processed_signals(processed_signals, t, dec=10, title='Full Signals', vert_marks=[],
                           filter_outliers=False, pos_outlier_zthresh=15, neg_outlier_zthresh=-5, t_min=0, t_max=np.inf):

    if utils.is_dict(list(processed_signals.values())[0]):
        n_panel_stacks = len(processed_signals.values())
    else:
        n_panel_stacks = 1
        # make a temporary outer dictionary for ease of use with for loop
        processed_signals = {'temp': processed_signals}

    t = t[::dec].copy()

    t_min_idx = np.argmax(t > t_min)
    t_max_idx = np.argwhere(t < t_max)[-1,0]

    # filter t but keep the same shape as signals
    t[:t_min_idx] = np.nan
    t[t_max_idx:] = np.nan

    # plot the raw signals and their baseline fits, baseline corrected signals, raw ligand and fitted iso, dff and baseline corrected df
    fig, axs = plt.subplots(2*n_panel_stacks, 2, layout='constrained', figsize=[20,6*n_panel_stacks], sharex=True)
    plt.suptitle(title)

    if len(vert_marks) > 0:
        vert_marks = vert_marks[(vert_marks > t_min) & (vert_marks < t_max)]

    for i, (sub_key, sub_signals) in enumerate(processed_signals.items()):

        # remove outliers in z-score space
        if filter_outliers:
            filt_sel = (sub_signals['z_dff_iso'][::dec] > neg_outlier_zthresh) & (sub_signals['z_dff_iso'][::dec] < pos_outlier_zthresh)
        else:
            filt_sel = np.full(t.shape, True)

        filt_t = t[filt_sel]

        gen_sub_title = sub_key + ' {}' if sub_key != 'temp' else '{}'

        # plot raw signals and baseline
        ax = axs[i,0]
        l3 = ax.plot(t, sub_signals['raw_iso'][::dec], label='Raw Iso', color='C1', alpha=0.5)
        l4 = ax.plot(t, sub_signals['baseline_iso'][::dec], '--', label='Iso Baseline', color='C1')
        l1 = ax.plot(t, sub_signals['raw_lig'][::dec], label='Raw Lig', color='C0', alpha=0.5)
        l2 = ax.plot(t, sub_signals['baseline_lig'][::dec], '--', label='Lig Baseline', color='C0')
        plot_utils.plot_dashlines(vert_marks, ax=ax)
        ax.set_title(gen_sub_title.format('Raw Signals'))
        ax.set_xlabel('Time (s)')
        ax.set_ylabel('Fluorescent Signal (V)')
        ax.xaxis.set_tick_params(which='both', labelbottom=True)
        ls = [l1[0], l2[0], l3[0], l4[0]]
        labs = [l.get_label() for l in ls]
        ax.legend(ls, labs, loc='center right')

        # plot baseline corrected signals
        ax = axs[i,1]
        ax.plot(filt_t, sub_signals['baseline_corr_lig'][::dec][filt_sel], label='Baseline Corrected Lig', alpha=0.5)
        ax.plot(filt_t, sub_signals['baseline_corr_iso'][::dec][filt_sel], label='Baseline Corrected Iso', alpha=0.5)
        ax.plot(filt_t, sub_signals['fitted_baseline_corr_iso'][::dec][filt_sel], label='Fitted Baseline Corrected Iso', alpha=0.5)
        plot_utils.plot_dashlines(vert_marks, ax=ax)
        ax.set_title(gen_sub_title.format('Baseline Subtracted Signals'))
        ax.set_xlabel('Time (s)')
        ax.set_ylabel('Fluorescent Signal (dV)')
        ax.xaxis.set_tick_params(which='both', labelbottom=True)
        ax.legend(loc='upper right')

        # plot raw ligand and fitted iso
        ax = axs[n_panel_stacks+i,0]
        ax.plot(filt_t, sub_signals['raw_lig'][::dec][filt_sel], label='Raw Lig', alpha=0.5)
        ax.plot(filt_t, sub_signals['fitted_iso'][::dec][filt_sel], label='Fitted Iso', alpha=0.5)
        plot_utils.plot_dashlines(vert_marks, ax=ax)
        ax.set_title(gen_sub_title.format('Iso ΔF/F Signals'))
        ax.set_xlabel('Time (s)')
        ax.set_ylabel('Fluorescent Signal (V)')
        ax.xaxis.set_tick_params(which='both', labelbottom=True)
        ax.legend(loc='upper right')

        # plot iso dFF and baseline corrected dFF
        ax = axs[n_panel_stacks+i,1]
        l2 = ax.plot(filt_t, sub_signals['dff_iso_baseline'][::dec][filt_sel], label='Baseline Corrected ΔF/F', color='C1', alpha=0.5)
        l1 = ax.plot(filt_t, sub_signals['dff_iso'][::dec][filt_sel], label='Iso ΔF/F', color='C0', alpha=0.5)
        plot_utils.plot_dashlines(vert_marks, ax=ax)
        ax.set_title(gen_sub_title.format('Iso Corrected Ligand Signals'))
        ax.set_xlabel('Time (s)')
        ax.set_ylabel('% ΔF/F')
        ax.xaxis.set_tick_params(which='both', labelbottom=True)
        ls = [l1[0], l2[0]]
        labs = [l.get_label() for l in ls]
        ax.legend(ls, labs, loc='upper right')

    return fig


def view_signal(processed_signals, t, signal_type, title=None, dec=10, vert_marks=[],
                filter_outliers=False, outlier_zthresh=10, t_min=0, t_max=np.inf, figsize=None,
                ylabel=''):

    if utils.is_dict(list(processed_signals.values())[0]):
        n_panel_stacks = len(processed_signals.values())
    else:
        n_panel_stacks = 1
        # make a temporary outer dictionary for ease of use with for loop
        processed_signals = {'temp': processed_signals}

    t = t[::dec]

    t_min_idx = np.argmax(t > t_min)
    t_max_idx = np.argwhere(t < t_max)[-1,0]

    if title is None:
        title = signal_type

    if figsize is None:
        figsize = (7, 3*n_panel_stacks)

    _, axs = plt.subplots(n_panel_stacks, 1, layout='constrained', figsize=figsize)
    plt.suptitle(title)

    if len(vert_marks) > 0:
        vert_marks = vert_marks[(vert_marks > t_min) & (vert_marks < t_max)]

    for i, (sub_key, sub_signals) in enumerate(processed_signals.items()):
        # remove outliers in z-score space
        filt_sel = np.full(t.shape, True)
        if filter_outliers:
            # filter based on raw signals
            z_sig = utils.z_score(sub_signals[signal_type][::dec])
            filt_sel = filt_sel & (np.abs(z_sig) < outlier_zthresh)

        # repurpose outlier filter for time filter as well
        filt_sel[:t_min_idx] = False
        filt_sel[t_max_idx:] = False

        filt_t = t[filt_sel]

        # plot raw signals and baseline
        ax = axs[i]
        ax.plot(filt_t, sub_signals[signal_type][::dec][filt_sel])
        plot_utils.plot_dashlines(vert_marks, ax=ax)
        ax.set_title(sub_key)
        ax.set_xlabel('Time (s)')
        ax.set_ylabel(ylabel)


def plot_aligned_signals(signal_dict, t, title, sub_titles_dict, x_label, y_label,
                         cmap = 'viridis', outlier_thresh=10, trial_markers=None):
    ''' Plot aligned signals in the given nested dictionary where the first set of keys define the rows and the
        second set of keys in the dictionaries indexed by the first set of keys define the columns

        In each 'cell' will plot a heatmap of the aligned signals above an average signal trace
    '''

    outer_keys = [k for k in signal_dict.keys() if k != 't']
    n_rows = len(outer_keys)
    n_cols = len(signal_dict[outer_keys[0]])

    fig, axs = plt.subplots(n_rows*2, n_cols, height_ratios=np.tile([3,2], n_rows),
                            figsize=(5*n_cols, 4*n_rows), layout='constrained')

    plt.suptitle(title)

    for i, key in enumerate(outer_keys):

        # remove outliers in z-score space
        z_signals = [utils.z_score(signal) for signal in signal_dict[key].values() if len(signal) > 0]
        signals_no_outliers = {}
        for j, (name, signal) in enumerate(signal_dict[key].items()):
            if len(signal) > 0:
                signal_no_outliers = signal.copy()
                signal_no_outliers[np.abs(z_signals[j]) > outlier_thresh] = np.nan
                signals_no_outliers[name] = signal_no_outliers
            else:
                signals_no_outliers[name] = signal

        max_act = np.max([np.max(signal[~np.isnan(signal)]) for signal in signals_no_outliers.values() if len(signal) > 0])
        min_act = np.min([np.min(signal[~np.isnan(signal)]) for signal in signals_no_outliers.values() if len(signal) > 0])

        # plot average signals first to get the x axis labels for the heatmap
        for j, (name, signal) in enumerate(signals_no_outliers.items()):
            hm_ax = axs[i*2, j]
            hm, _ = plot_utils.plot_stacked_heatmap_avg(signal, t, hm_ax, axs[i*2+1, j],
                                             x_label=x_label, y_label=y_label,
                                             title='{} {}'.format(key, sub_titles_dict[name]),
                                             show_cbar=False, cmap=cmap, vmax=max_act, vmin=min_act)

            # plot any horizontal trial markers
            if not trial_markers is None:
                if name in trial_markers:
                    for t_idx in trial_markers[name]:
                        hm_ax.axhline(t_idx - 0.5, dashes=[4, 4], c='k', lw=1)

        # share y axis for all average rate plots
        # find largest y range
        min_y = np.min([ax.get_ylim()[0] for ax in axs[i*2+1, :]])
        max_y = np.max([ax.get_ylim()[1] for ax in axs[i*2+1, :]])
        for j in range(axs.shape[1]):
            axs[i*2+1, j].set_ylim(min_y, max_y)

        # create color bar legend for heatmap plots
        fig.colorbar(hm, ax=axs[i*2,:].ravel().tolist(), label=y_label)


def plot_avg_signals(plot_groups, group_labels, data_mat_dict, regions, t, title, plot_titles, x_label, y_label,
                     xlims_dict=None, dashlines=None, legend_params=None, use_se=True, group_colors=None, ph=3.5, pw=5):

    def calc_error(mat):
        if use_se:
            return utils.stderr(mat, axis=0)
        else:
            return np.std(mat, axis=0)

    plotted = False

    n_rows = len(regions)
    n_cols = len(plot_groups)
    fig, axs = plt.subplots(n_rows, n_cols, figsize=(pw*n_cols, ph*n_rows), layout='constrained', sharey='row')

    if n_rows == 1:
        axs = axs[None,:]
    if n_cols == 1:
        axs = axs[:,None]

    fig.suptitle(title)

    for i, region in enumerate(regions):

        # limit x axis like this so that the y is scaled to what is plotted
        if xlims_dict is None:
            t_sel = np.full(t.shape, True)
        else:
            t_sel = (t > xlims_dict[region][0]) & (t < xlims_dict[region][1])

        for j, (plot_group, plot_title) in enumerate(zip(plot_groups, plot_titles)):
            ax = axs[i,j]
            ax.set_title(region + ' ' + plot_title)
            if not dashlines is None:
                plot_utils.plot_dashlines(dashlines, ax=ax)

            for k, group in enumerate(plot_group):
                if group in data_mat_dict[region]:
                    act = data_mat_dict[region][group]
                    if len(act) > 0:
                        plotted = True
                        if not group_colors is None:
                            plot_utils.plot_psth(t[t_sel], np.nanmean(act, axis=0)[t_sel], calc_error(act)[t_sel], ax, label=group_labels[group], color=group_colors[k])
                        else:
                            plot_utils.plot_psth(t[t_sel], np.nanmean(act, axis=0)[t_sel], calc_error(act)[t_sel], ax, label=group_labels[group])

            if j == 0:
                ax.set_ylabel(y_label)
            else:
                ax.yaxis.set_tick_params(which='both', labelleft=True)
            ax.set_xlabel(x_label)

            if not legend_params is None:
                if not legend_params[region] is None:
                    ax.legend(**legend_params[region])
            else:
                ax.legend(loc='best')

    return fig, plotted


def get_signal_type_labels(signal_type):
    ''' Get signal titles and labels based on the type of signal '''
    match signal_type:
        case 'raw_lig':
            title = 'Raw Ligand Signal'
            ax_label = 'Signal (V)'
        case 'raw_iso':
            title = 'Raw Isosbestic Signal'
            ax_label = 'Signal (V)'
        case 'filtered_lig':
            title = 'Filtered Ligand Signal'
            ax_label = 'Signal (V)'
        case 'filtered_iso':
            title = 'Filtered Isosbestic Signal'
            ax_label = 'Signal (V)'
        case 'baseline_corr_lig':
            title = 'Baseline-Subtracted Ligand Signal'
            ax_label = 'ΔF'
        case 'baseline_corr_iso':
            title = 'Baseline-Subtracted Isosbestic Signal'
            ax_label = 'ΔF'
        case 'fitted_iso':
            title = 'Fitted Isosbestic Signal'
            ax_label = 'Signal (V)'
        case 'fitted_baseline_corr_iso':
            title = 'Fitted Baseline-Subtracted Isosbestic Signal'
            ax_label = 'ΔF'
        case 'dff_iso':
            title = 'Isosbestic Normalized Signal'
            ax_label = '% dF/F'
        case 'z_dff_iso':
            title = 'Z-scored Isosbestic Normalized Signal'
            ax_label = 'Z-scored dF/F'
        case 'dff_iso_baseline':
            title = 'Baseline-Subtracted Isosbestic Normalized Signal'
            ax_label = '% dF/F'
        case 'z_dff_iso_baseline':
            title = 'Z-scored Baseline-Subtracted Isosbestic Normalized Signal'
            ax_label = 'Z-scored dF/F'


    return title, ax_label


def get_implant_side_type(side, implant_side):
    return 'ipsi' if implant_side == side else 'contra'


def plot_power_spectra(signals, dt, f_max=20, title='', log_x=True, signal_names=None):
    #signal = signal - np.mean(signal)
    # ps = np.abs(np.fft.rfft(signal))**2
    # freqs = np.fft.rfftfreq(signal.size, dt)
    n_sigs = len(signals)

    if signal_names is None:
        signal_names = ['' for i in range(n_sigs)]

    fig, ax = plt.subplots(1)

    if log_x:
        x_lims = [0.001, f_max]
    else:
        x_lims = [0, f_max]

    for i in range(n_sigs):
        signal = signals[i]

        # get rid of nans
        tmp_sig = signal[~np.isnan(signal)]

        freqs, ps = sig.welch(tmp_sig, fs = 1/dt, nperseg = round(1/dt)*30, scaling='density')

        #freqs, ps = sig.periodogram(signal, fs = 1/dt, scaling='spectrum')
        # same as FFT but simpler

        freq_sel = (freqs > x_lims[0]) & (freqs < x_lims[1])

        ax.plot(freqs[freq_sel], ps[freq_sel], label=signal_names[i])

    ax.set_yscale('log')
    if log_x:
        ax.set_xscale('log')

    ax.set_xlabel('Frequency (Hz)')
    ax.set_ylabel('Power Spectral Density (V^2/Hz)')
    ax.set_title(title)
    ax.legend()

    return fig


def remove_outliers(mat, outlier_thresh):
    mat[np.abs(mat) > outlier_thresh] = np.nan

    return mat


def stack_fp_mats(mat_dict, regions, sess_ids, subjects, signal_type, filter_outliers=False, outlier_thresh=20, groups=None):

    if not utils.is_list(subjects):
        subjects = [subjects]

    stacked_mats = {region: {} for region in regions}
    if groups is None:
        groups = mat_dict[sess_ids[subjects[0]][0]][signal_type][regions[0]].keys()

    for region in regions:
        for group in groups:
            for subj_id in subjects:
                mat = np.vstack([mat_dict[sess_id][signal_type][region][group] for sess_id in sess_ids[subj_id] if group in mat_dict[sess_id][signal_type][region].keys()])
                if filter_outliers:
                    mat = remove_outliers(mat, outlier_thresh)

                # TODO: This does not work for multiple subjects
                stacked_mats[region][group] = mat

    return stacked_mats


# def combine_group_mats(data_dict, grouping=None):
#     if grouping is None:
#         grouping = {'all': list(data_dict[regions[0]].keys())}

#     for region in regions:
#         for name, groups in grouping.items():
#             data_dict[region][name] = np.vstack([data_dict[region][group] for group in groups])

#     return data_dict




# %% Alignment Helpers

class Alignment(StrEnum):
    cport_on = 'cport_on'
    cpoke_in = 'cpoke_in'
    early_cpoke_in = 'early_cpoke_in'
    tone = 'tone'
    cue = 'cue'
    cpoke_out = 'cpoke_out'
    early_cpoke_out = 'early_cpoke_out'
    resp = 'resp'
    reward = 'reward'
    cue_poke_resp = 'norm_cue_poke_resp'
    poke_cue_resp = 'norm_poke_cue_resp'
    resp_reward = 'norm_resp_reward'
    cue_resp = 'norm_cue_resp'

def get_align_title(align):
    match align:
        case Alignment.cport_on:
            return 'Center Port On'
        case Alignment.cpoke_in:
            return 'Center Poke In'
        case Alignment.early_cpoke_in:
            return 'Early Center Poke In'
        case Alignment.tone:
            return 'Tone Start'
        case Alignment.cue:
            return 'Response Cue'
        case Alignment.cpoke_out:
            return 'Center Poke Out'
        case Alignment.early_cpoke_out:
            return 'Early Center Poke Out'
        case Alignment.resp:
            return 'Response Poke'
        case Alignment.reward:
            return 'Reward Delivery'
        case Alignment.cue_poke_resp:
            return 'Normalized Cue, Poke Out, & Response'
        case Alignment.poke_cue_resp:
            return 'Normalized Poke Out, Cue, & Response'
        case Alignment.resp_reward:
            return 'Normalized Response to Reward'
        case Alignment.cue_resp:
            return 'Normalized Cue to Response'
        case _:
            return align

def get_align_xlabel(align):
    gen_x_label = 'Time from {} (s)'
    gen_norm_x_label = 'Normalized Time from {}'
    match align:
        case Alignment.cport_on:
            return gen_x_label.format('port on')
        case Alignment.cpoke_in:
            return gen_x_label.format('poke in')
        case Alignment.early_cpoke_in:
            return gen_x_label.format('poke in')
        case Alignment.tone:
            return gen_x_label.format('tone onset')
        case Alignment.cue:
            return gen_x_label.format('response cue')
        case Alignment.cpoke_out:
            return gen_x_label.format('poke out')
        case Alignment.early_cpoke_out:
            return gen_x_label.format('poke out')
        case Alignment.resp:
            return gen_x_label.format('response poke')
        case Alignment.reward:
            return gen_x_label.format('reward')
        case Alignment.cue_poke_resp:
            return gen_norm_x_label.format('response cue to response')
        case Alignment.poke_cue_resp:
            return gen_norm_x_label.format('poke out to response')
        case Alignment.resp_reward:
            return gen_norm_x_label.format('response to reward')
        case Alignment.cue_resp:
            return gen_norm_x_label.format('cue to response')


# %% Plot Image Saving/Loading Methods

def save_fig(fig, save_path, format='png', **kwargs):

    utils.check_make_dir(save_path)

    if not format in save_path:
        save_path += '.' + format
        
    if format == 'svg':
        plt.rcParams["svg.fonttype"] = 'none'

    fig.savefig(save_path, **kwargs)


def get_base_figure_save_path():
    # make sure subject id is a string
    return path.join(utils.get_user_home(), 'FP Images')


def get_figure_save_path(behavior, subj_ids, filename=''):
    # make sure subject id is a string
    subj_ids = str(subj_ids) if not utils.is_list(subj_ids) else '_'.join([str(i) for i in sorted(subj_ids)])
    return path.join(get_base_figure_save_path(), behavior, subj_ids, filename)


def generate_figure_ppt(save_path, group_by=None, behaviors=None, subjects=None, alignments=None, filenames=None):
    '''
    Generate a power point from saved figures for the specified behaviors, subjects, and alignments grouped by the provided ordering

    Parameters
    ----------
    save_path : The path to the saved powerpoint
    group_by : The grouping order for the images. The default is ['behavior', 'subject', 'alignment', 'filename'].
        This determines the order in which images are placed in the powerpoint
    behaviors : The behaviors to include in the powerpoint. The default is all behaviors.
    subjects : The subject ids to include in the powerpoint. The default is all subjects.
    alignments : The alignment points to include in the powerpoint. The default is all alignments.

    Returns
    -------
    None.

    '''

    base_path = get_base_figure_save_path()
    def_group_by = ['behavior', 'subject', 'alignment', 'filename']
    if group_by is None:
        group_by = def_group_by
    elif not utils.is_list(group_by):
        group_by = [group_by]

    if len(group_by) < len(def_group_by):
        group_by.extend([g for g in def_group_by if g not in group_by])

    if behaviors is None:
        behaviors = [f.name for f in os.scandir(base_path) if f.is_dir()]
    elif not utils.is_list(behaviors):
        # make sure a single behavior is in list form
        behaviors = [behaviors]

    if subjects is None:
        subj_id_paths = utils.flatten([glob(path.join(base_path, beh, '*')) for beh in behaviors])
        subjects = np.unique([path.basename(p) for p in subj_id_paths]).tolist()
    elif not utils.is_list(subjects):
        subjects = [subjects]

    # make sure subject ids are strings and collapse multiple ids into
    subjects = [str(i) if not utils.is_list(i) else '_'.join([str(j) for j in sorted(i)]) for i in subjects]

    if alignments is None:
        alignments = [a for a in Alignment]
    elif not utils.is_list(alignments):
        # make sure a single alignment is in list form
        alignments = [alignments]

    if filenames is None:
        # get all possible filenames of figures
        all_filenames = utils.flatten([glob(path.join(base_path, b, s, a+'*')) for b in behaviors for s in subjects for a in alignments])
        # remove path and extension
        filenames = np.unique([path.basename(f).replace('.png', '') for f in all_filenames])
        # remove alignment
        filenames = [f.replace(a+'_', '', 1) for a in alignments for f in filenames if a in f]
    elif not utils.is_list(filenames):
        # make sure a single filename is in list form
        filenames = [filenames]

    # make sure there are no duplicates and the order is preserved if passed in
    _, idx = np.unique(filenames, return_index=True)
    filenames = np.array(filenames)[np.sort(idx)]

    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    __add_ppt_slides_recursive(prs, group_by, behaviors, subjects, alignments, filenames, 0)

    if not '.pptx' in save_path:
        save_path += '.pptx'
    prs.save(save_path)


def __add_ppt_slides_recursive(prs, group_by, behaviors, subjects, alignments, filenames, level, level_titles=[]):

    group = group_by[level]
    title_font_size = Pt(60 - 10*level)

    match group:
        case 'behavior':
            level_groups = behaviors
            level_group_labels = behaviors
        case 'subject':
            level_groups = subjects
            level_group_labels = [s.replace('_', ', ') for s in subjects]
        case 'alignment':
            level_groups = alignments
            level_group_labels = [get_align_title(a) for a in alignments]
        case 'filename':
            level_groups = filenames
            level_group_labels = filenames

    # if we haven't reached the lowest level make a title slide for each level group and recurse
    if level < len(group_by)-1:
        for i, level_group in enumerate(level_groups):

            match group:
                case 'behavior':
                    images_exist = __check_images_exist(level_group, subjects, alignments, filenames)
                    need_title = __check_images_exist(level_group, subjects, alignments, filenames, count_thresh=1)
                case 'subject':
                    images_exist = __check_images_exist(behaviors, level_group, alignments, filenames)
                    need_title = __check_images_exist(behaviors, level_group, alignments, filenames, count_thresh=1)
                case 'alignment':
                    images_exist = __check_images_exist(behaviors, subjects, level_group, filenames)
                    need_title = __check_images_exist(behaviors, subjects, level_group, filenames, count_thresh=1)
                case 'filename':
                    images_exist = __check_images_exist(behaviors, subjects, alignments, level_group)
                    need_title = __check_images_exist(behaviors, subjects, alignments, level_group, count_thresh=1)

            if images_exist:
                level_titles_copy = level_titles.copy()
                level_titles_copy.append(level_group_labels[i])

                if need_title:
                    layout = prs.slide_layouts[0] # title slide
                    slide = prs.slides.add_slide(layout)
                    title = slide.shapes.title
                    title.text = '\n'.join(level_titles_copy)
                    for line in title.text_frame.paragraphs:
                        line.font.size = title_font_size
                        if level == 0:
                            line.font.bold = True

                    # need to reset the height and top after changing width and left
                    title_height = title.height
                    title_top = title.top
                    title.width = prs.slide_width
                    title.left = 0
                    title.height = title_height
                    title.top = title_top

                match group:
                    case 'behavior':
                        __add_ppt_slides_recursive(prs, group_by, level_group, subjects, alignments, filenames, level+1, level_titles_copy)
                    case 'subject':
                        __add_ppt_slides_recursive(prs, group_by, behaviors, level_group, alignments, filenames, level+1, level_titles_copy)
                    case 'alignment':
                        __add_ppt_slides_recursive(prs, group_by, behaviors, subjects, level_group, filenames, level+1, level_titles_copy)
                    case 'filename':
                        __add_ppt_slides_recursive(prs, group_by, behaviors, subjects, alignments, level_group, level+1, level_titles_copy)

    else:
        # we are at the lowest level so we need to start adding slides
        for i, level_group in enumerate(level_groups):
            # get all images the for the given groups
            # Note: at the lowest level there will be one element for each of the grouping arguments
            match group:
                case 'behavior':
                    image_files = glob(get_figure_save_path(level_group, subjects, alignments+'*'+filenames+'.*'))
                case 'subject':
                    image_files = glob(get_figure_save_path(behaviors, level_group, alignments+'*'+filenames+'.*'))
                case 'alignment':
                    image_files = glob(get_figure_save_path(behaviors, subjects, level_group+'*'+filenames+'.*'))
                case 'filename':
                    image_files = glob(get_figure_save_path(behaviors, subjects, alignments+'*'+level_group+'.*'))

            # TODO: Only match exact filenames, not any pattern match
            # The tradeoff is having the * in the search to match files without an alignment. So perhaps just switch on whether there is an alignment or not
            image_files = sorted(image_files)
            level_titles_copy = level_titles.copy()
            level_titles_copy.append(level_group_labels[i])

            if len(image_files) > 1:
                layout = prs.slide_layouts[0] # title slide
                slide = prs.slides.add_slide(layout)
                title = slide.shapes.title
                title.text = '\n'.join(level_titles_copy)
                for line in title.text_frame.paragraphs:
                    line.font.size = title_font_size

                # need to reset the height and top after changing width and left
                title_height = title.height
                title_top = title.top
                title.width = prs.slide_width
                title.left = 0
                title.height = title_height
                title.top = title_top

            for image_file in image_files:

                layout = prs.slide_layouts[5] # title only
                slide = prs.slides.add_slide(layout)
                title = slide.shapes.title
                title.text = ', '.join(level_titles_copy)
                title.text_frame.paragraphs[0].font.size = title_font_size
                # need to reset the height and top after changing width and left
                title_height = title.height
                title.width = prs.slide_width
                title.height = title_height
                title.left = 0
                title.top = 0

                image = slide.shapes.add_picture(image_file, 0, 0, width=prs.slide_width - Inches(0.1))

                # rescale image if too tall
                if image.height > prs.slide_height - title.height:
                    aspect = image.height/image.width
                    image.height = prs.slide_height - title.height - Inches(0.1)
                    image.width = int(image.height/aspect)

                image.left = int((prs.slide_width - image.width) / 2)
                image.top = int((prs.slide_height - image.height - title.height) / 2) + title.height

def __check_images_exist(behaviors, subjects, alignments, filenames, count_thresh = 0):

    if not utils.is_list(behaviors):
        behaviors = [behaviors]

    if not utils.is_list(subjects):
        subjects = [subjects]

    if not utils.is_list(alignments):
        alignments = [alignments]

    if not utils.is_list(filenames):
        filenames = [filenames]

    file_count = 0

    for b in behaviors:
        for s in subjects:
            for a in alignments:
                for f in filenames:
                    files = glob(path.join(get_base_figure_save_path(), b, s, a+'*'+f+'.*'))
                    if len(files) > 0:
                        file_count += 1
                    if file_count > count_thresh:
                        return True
                if file_count > count_thresh:
                    return True
            if file_count > count_thresh:
                return True
        if file_count > count_thresh:
            return True

    return False
